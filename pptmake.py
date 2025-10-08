from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation
prs = Presentation()

# Define slide layout
title_slide_layout = prs.slide_layouts[0]
bullet_slide_layout = prs.slide_layouts[1]

# Helper function to add bullet slide
def add_bullet_slide(title, bullets):
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = title
    tf = body_shape.text_frame
    tf.clear()
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)

# 1. Title Slide
slide = prs.slides.add_slide(title_slide_layout)
slide.shapes.title.text = "Elevating IT Quality in Banking"
slide.placeholders[1].text = "Strategic Assessment and Integration Approaches"

# 2. Strategic Challenge
add_bullet_slide("Strategic Challenge: Need for IT Quality Assessment", [
    "Increasing reliance on IT systems in banking operations",
    "Customer expectations for seamless digital experiences",
    "Regulatory pressures and compliance requirements",
    "Need for operational efficiency and risk mitigation"
])

# 3. Overview of Approaches
add_bullet_slide("Overview of Key Approaches", [
    "Customer Satisfaction Surveys",
    "Service Level Agreements (SLAs)",
    "Maturity Models & Norm Compliance",
    "SERVQUAL",
    "ISO/IEC 20000",
    "COBIT 5"
])

# 4. Detailed Slides for Each Approach
approaches = {
    "Customer Satisfaction Surveys": [
        "Purpose: Measure user perception of IT services",
        "Example: Post-service feedback forms",
        "Benefits: Direct insights, continuous improvement",
        "Risks: Biased responses, low participation",
        "Mitigation: Anonymity, incentives"
    ],
    "Service Level Agreements (SLAs)": [
        "Purpose: Define expected service performance",
        "Example: 99.9% uptime commitment",
        "Benefits: Accountability, clear expectations",
        "Risks: Misalignment, unrealistic targets",
        "Mitigation: Regular reviews, stakeholder input"
    ],
    "Maturity Models & Norm Compliance": [
        "Purpose: Assess IT process maturity",
        "Example: CMMI, ITIL maturity levels",
        "Benefits: Benchmarking, structured growth",
        "Risks: Complexity, resource demands",
        "Mitigation: Phased implementation, training"
    ],
    "SERVQUAL": [
        "Purpose: Evaluate service quality gaps",
        "Example: Tangibles, reliability, responsiveness",
        "Benefits: Holistic view, customer-centric",
        "Risks: Subjectivity, cultural bias",
        "Mitigation: Standardized surveys, diverse sampling"
    ],
    "ISO/IEC 20000": [
        "Purpose: Certify IT service management quality",
        "Example: Compliance audits",
        "Benefits: Global recognition, process discipline",
        "Risks: Audit fatigue, cost",
        "Mitigation: Internal audits, budget planning"
    ],
    "COBIT 5": [
        "Purpose: Governance and management framework",
        "Example: Process capability assessments",
        "Benefits: Strategic alignment, risk control",
        "Risks: Overhead, complexity",
        "Mitigation: Tailored adoption, expert guidance"
    ]
}

for title, bullets in approaches.items():
    add_bullet_slide(title, bullets)

# 5. Integration Strategy
add_bullet_slide("Integration Strategy", [
    "Combine quantitative and qualitative assessments",
    "Align SLAs with customer feedback",
    "Use maturity models to guide ISO/COBIT adoption",
    "Cross-reference SERVQUAL with survey data",
    "Create unified dashboard for IT quality metrics"
])

# 6. Real-World Success Stories
slide = prs.slides.add_slide(bullet_slide_layout)
slide.shapes.title.text = "Real-World Success Stories"
rows = [
    ["Bank", "Approaches Used", "Outcomes"],
    ["NordBank", "SLAs, ISO/IEC 20000", "Improved uptime, audit success"],
    ["BalticTrust", "Surveys, SERVQUAL", "Higher customer satisfaction"],
    ["EstoniaBank", "COBIT 5, Maturity Models", "Strategic IT alignment"]
]
text = "\n".join([" | ".join(row) for row in rows])
slide.placeholders[1].text = text

# 7. Final Outcomes
add_bullet_slide("Final Outcomes: Strategic Benefits Achieved", [
    "Enhanced customer trust and loyalty",
    "Improved regulatory compliance",
    "Optimized IT resource allocation",
    "Increased operational resilience"
])

# 8. Closing Message
add_bullet_slide("Closing Message", [
    "\"Quality in IT is not a destination, but a journey of continuous excellence.\"",
    "— CEO, Global Banking Group"
])

# Save presentation
prs.save("/mnt/data/IT_Quality_Assessment_Banking.pptx")
print("Presentation created successfully.")
